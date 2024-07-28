import streamlit as st
from pptx import Presentation
import os
import google.generativeai as genai

# Configure the Generative AI API
api = "AIzaSyDFZPjLcX0l__4Xl9ZuWX7bKNr6a3CvIQo"
genai.configure(api_key=api)

# Create the model
generation_config = {
    "temperature": 1,
    "top_p": 0.95,
    "top_k": 64,
    "max_output_tokens": 8192,
}

model = genai.GenerativeModel(
    model_name="gemini-1.5-pro",
    generation_config=generation_config,
)

chat_session = model.start_chat(
    history=[
        {
            "role": "user",
            "parts": [
                "I need your help generating content for my PowerPoint presentation. Please follow this format:\n\nEach slide title should be placed before the colon (:)\nEach slide content should be placed after the equals sign (=)\nEach new slide should start on a new line\nFor example:\n\nTitle1: Content for the first slide =\nTitle2: Content for the second slide =\nTitle3: Content for the third slide =\n\nPlease ensure that the content is relevant to the given titles. Thank you!",
            ],
        },
        {
            "role": "model",
            "parts": [
                "Please provide me with the slide titles and any specific content ideas you have for each slide. For example:\n\n* **Slide 1: Introduction** = A brief overview of the presentation topic.\n* **Slide 2: Problem Statement** = Description of the problem being addressed. \n* **Slide 3: Proposed Solution** =  Details about the proposed solution to the problem.\n\nThe more information you give me about your presentation, the better I can tailor the content to your specific needs.",
            ],
        },
        {
            "role": "user",
            "parts": [
                "I need your help transforming user input into a specific format for my PowerPoint presentation. The user will enter the information in the following way:\nTopic: [topic]\nPage 1: [title]\nPage 2: [title]\nPage 3: [title]\n...\nPlease transform this input into the following format:\ntitle1: content1 =\ntitle2: content2 =\ntitle3: content3 =\n...\nFor example:\nUser Input:\nTopic: Fruit\nPage 1: First Benefit\nPage 2: Second Benefit\nPage 3: Content\n\nFirst Benefit: content1 =\nSecond Benefit: content2 =\nContent: content3 =",
            ],
        },
        {
            "role": "model",
            "parts": [
                "```python\ndef format_presentation_input(user_input):\n  \"\"\"Transforms user input into a PowerPoint slide format.\n\n  Args:\n    user_input: A string containing the user's input in the format:\n      \"Topic: [topic]\\nPage 1: [title]\\nPage 2: [title]\\n...\".\n\n  Returns:\n    A string representing the formatted output for PowerPoint slides, \n    or an error message if the input is invalid.\n  \"\"\"\n\n  lines = user_input.strip().split('\\n')\n  if len(lines) < 2 or lines[0].startswith(\"Page\"):\n    return \"Invalid input format. Please start with 'Topic: [topic]'.\"\n\n  topic = lines[0].split(\":\")[1].strip()\n  formatted_output = \"\"\n\n  for line in lines[1:]:\n    if \":\" not in line:\n      return \"Invalid input format for page titles.\"\n    page, title = line.split(\":\", 1)\n    title = title.strip()\n    formatted_output += f\"{title}: content = \\n\"  \n\n  return formatted_output\n\n# Example usage:\nuser_input = \"\"\"Topic: Fruit\nPage 1: First Benefit\nPage 2: Second Benefit\nPage 3: Content\"\"\"\n\nformatted_presentation = format_presentation_input(user_input)\nprint(formatted_presentation)\n```\n\n**Explanation:**\n\n1. **Function Definition:** Defines a function `format_presentation_input()` that takes the user's input as a single string.\n2. **Input Validation:** Checks for minimum input length and ensures the first line starts with \"Topic:\". If invalid, returns an error message.\n3. **Extracting Topic:** Extracts the topic from the first line of the input.\n4. **Iterating and Formatting:** Iterates through the remaining lines (representing pages) and extracts the title from each line. \n5. **Constructing Output:**  Formats each extracted title into the desired `title: content =` format and appends it to the `formatted_output` string. \n6. **Returning Output:** Returns the `formatted_output` string, which contains the formatted slide titles.\n\nNow, you can use this Python code to process the user input and generate the desired format for your PowerPoint presentation.",
            ],
        },
        {
            "role": "user",
            "parts": [
                "Topic: Sport\nPage 1: First Benefit\nPage 2: Second Benefit\nPage 3: Example",
            ],
        },
        {
            "role": "model",
            "parts": [
                "Physical Prowess and Well-being: Engaging in sports provides numerous physical benefits, boosting cardiovascular health, strength, coordination, and overall fitness. It encourages an active lifestyle and combats sedentary behavior, promoting long-term well-being. (Visual: Action shot of athletes showcasing athleticism, infographic on health benefits=\n"
                "Mental Fortitude and Resilience:Sport demands discipline, focus, and determination. Athletes develop mental toughness, learning to overcome challenges, cope with pressure, and bounce back from setbacks. These skills extend far beyond the playing field, translating into success in various aspects of life. (Visual: Close-up of an athlete's focused expression, image representing perseverance=\n"
                "Teamwork and Leadership:Many sports rely on teamwork, fostering collaboration, communication, "
                "and leadership skills.=",
            ],
        },
    ]
)


def choose_ppt_style():
    styles = ["Simple and Minimalist", "Colorful and Vibrant", "Modern and Trendy"]
    style_choice = st.selectbox("Choose a PowerPoint design style:", styles)
    style_map = {
        "Simple and Minimalist": "theme/simple.pptx",
        "Colorful and Vibrant": "theme/colorful.pptx",
        "Modern and Trendy": "theme/facet.pptx"
    }
    return style_map[style_choice]


def get_user_input():
    user_title = st.text_input("Presentation Title:")
    topic = st.text_input("Topic:")
    num_pages = st.number_input("Number of Pages:", min_value=1, step=1)

    pages = []
    for x in range(1, num_pages + 1):
        page_content = st.text_input(f"Page {x}: Title", key=f"page_{x}")
        pages.append(f"Page {x}: {page_content}")

    style = choose_ppt_style()

    if st.button("Generate Presentation"):
        user_input = f"Topic: {topic}\n" + "\n".join(pages)
        return user_title, user_input, style

    return None, None, None


def format_presentation_input(user_input):
    lines = user_input.strip().split('\n')
    if len(lines) < 2 or not lines[0].startswith("Topic:"):
        return "Invalid input format. Please start with 'Topic: [topic]'."

    topic = lines[0].split(":")[1].strip()
    formatted_output = ""

    for line in lines[1:]:
        if ":" not in line:
            return "Invalid input format for page titles."
        page, title = line.split(":", 1)
        title = title.strip()
        formatted_output += f"{title}: content = \n"

    return formatted_output


def read_and_parse_file(file_path):
    titles = []
    contents = []
    pictures = []

    with open(file_path, 'r') as file:
        lines = file.readlines()

        for line in lines:
            if ':' in line:
                title_part, rest = line.split(':', 1)
                title = title_part.strip()
                titles.append(title)

                if '(' in rest:
                    content_part, picture_part = rest.split('(', 1)
                    content = content_part.strip()
                    picture = '(' + picture_part.strip()
                else:
                    content = rest.strip()
                    picture = ''

                contents.append(content)
                pictures.append(picture)

    return titles, contents, pictures


def create_presentation(user_title, titles, contents, style):
    prs = Presentation(style)

    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = user_title
    subtitle.text = "Generated Presentation"

    # Add slides for each title and content
    bullet_slide_layout = prs.slide_layouts[1]

    for title, content in zip(titles, contents):
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title

        try:
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.text = content
        except IndexError:
            st.warning(f"Placeholder not found for slide '{title}'. Content may be missing.")

        title_shape.text = title

    # Save the presentation
    prs.save(f"{user_title}.pptx")
    st.success(f"Presentation '{user_title}.pptx' created successfully.")
    with open(f"{user_title}.pptx", "rb") as file:
        st.download_button(
            label="Download Presentation",
            data=file,
            file_name=f"{user_title}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )


def main():
    st.title("PowerPoint Presentation Generator")

    user_title, user_input, style = get_user_input()

    if user_title and user_input:
        input_to_gemini = format_presentation_input(user_input)
        response = chat_session.send_message(input_to_gemini)
        with open("content", "w") as f:
            f.write(response.text)

        titles, contents, _ = read_and_parse_file("content")
        create_presentation(user_title, titles, contents, style)


if __name__ == "__main__":
    main()
